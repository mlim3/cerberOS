#!/usr/bin/env python3
"""Generate Aegis DataBus presentation PowerPoint."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util
from pptx.oxml.ns import qn
from pptx.enum.dml import MSO_THEME_COLOR
import copy
from lxml import etree

# ── Color Palette (deep navy + steel blue + white + accent gold) ──────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)   # dark background / header
STEEL      = RGBColor(0x1B, 0x4F, 0x72)   # secondary blue
ACCENT     = RGBColor(0xF0, 0xA5, 0x00)   # gold accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xEC, 0xF0, 0xF1)
DARK_GRAY  = RGBColor(0x2C, 0x3E, 0x50)
MID_BLUE   = RGBColor(0x21, 0x8B, 0xC3)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────
def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_layout(prs):
    return prs.slide_layouts[6]  # completely blank


def add_slide(prs):
    return prs.slides.add_slide(blank_layout(prs))


def rect(slide, l, t, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(1, l, t, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def txbox(slide, text, l, t, w, h,
          size=18, bold=False, color=WHITE,
          align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return tb


def header_bar(slide, title, subtitle=None):
    """Navy top bar with gold-underlined title."""
    rect(slide, 0, 0, SLIDE_W, Inches(1.35), fill=NAVY)
    # thin gold rule
    rect(slide, 0, Inches(1.35), SLIDE_W, Pt(4), fill=ACCENT)

    txbox(slide, title,
          Inches(0.45), Inches(0.15), Inches(12), Inches(0.7),
          size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txbox(slide, subtitle,
              Inches(0.45), Inches(0.85), Inches(12), Inches(0.45),
              size=14, color=ACCENT, align=PP_ALIGN.LEFT)


def body_bg(slide):
    rect(slide, 0, Inches(1.39), SLIDE_W, SLIDE_H - Inches(1.39), fill=LIGHT_GRAY)


def bullet_block(slide, items, l=Inches(0.5), t=Inches(1.65), w=Inches(12.3), item_h=Inches(0.48),
                 size=17, color=DARK_GRAY, bullet="▸ "):
    for i, item in enumerate(items):
        txbox(slide, f"{bullet}{item}", l, t + i * item_h, w, item_h,
              size=size, color=color)


def two_col_table(slide, headers, rows, l, t, w, h,
                  hdr_fill=NAVY, row_fill=WHITE, alt_fill=RGBColor(0xD6,0xEA,0xF8)):
    """Draw a simple table using rectangles + textboxes."""
    n_cols = len(headers)
    col_w = w // n_cols
    row_h = h // (len(rows) + 1)

    # Header row
    for ci, hdr in enumerate(headers):
        rect(slide, l + ci * col_w, t, col_w, row_h, fill=hdr_fill)
        txbox(slide, hdr, l + ci * col_w + Inches(0.1), t + Pt(4),
              col_w - Inches(0.15), row_h,
              size=13, bold=True, color=WHITE)

    for ri, row in enumerate(rows):
        fill = row_fill if ri % 2 == 0 else alt_fill
        for ci, cell in enumerate(row):
            rect(slide, l + ci * col_w, t + (ri + 1) * row_h, col_w, row_h,
                 fill=fill, line=RGBColor(0xBD,0xC3,0xC7))
            txbox(slide, cell, l + ci * col_w + Inches(0.08),
                  t + (ri + 1) * row_h + Pt(3),
                  col_w - Inches(0.1), row_h,
                  size=12, color=DARK_GRAY)


def add_speaker_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def add_footer(slide, text="CS-686 | Aegis DataBus | Meron Shibiru | March 9, 2026"):
    rect(slide, 0, SLIDE_H - Inches(0.32), SLIDE_W, Inches(0.32), fill=NAVY)
    txbox(slide, text, Inches(0.3), SLIDE_H - Inches(0.30),
          SLIDE_W - Inches(0.6), Inches(0.28),
          size=9, color=RGBColor(0xAE,0xBF,0xC8), align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    slide = add_slide(prs)
    # Full navy background
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    # Gold accent stripe left edge
    rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill=ACCENT)
    # Gold horizontal rule
    rect(slide, Inches(0.3), Inches(4.0), Inches(12.7), Pt(3), fill=ACCENT)

    # Main title
    txbox(slide, "Aegis DataBus",
          Inches(0.5), Inches(1.2), Inches(12.5), Inches(1.1),
          size=52, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    txbox(slide, "Central Message Transport for Aegis OS",
          Inches(0.5), Inches(2.2), Inches(12.5), Inches(0.65),
          size=24, bold=False, color=ACCENT, align=PP_ALIGN.LEFT)
    txbox(slide, "CS-686  ·  March 9, 2026  ·  Meron Shibiru",
          Inches(0.5), Inches(2.85), Inches(10), Inches(0.45),
          size=16, color=RGBColor(0xAE,0xBF,0xC8), align=PP_ALIGN.LEFT)

    # ── Simple ASCII-style architecture diagram (text-art boxes) ─────────────
    diag_l = Inches(0.6)
    diag_t = Inches(4.2)
    box_w  = Inches(1.55)
    box_h  = Inches(0.6)
    gap    = Inches(0.18)

    components = [
        ("I/O Layer",      MID_BLUE),
        ("Orchestrator",   STEEL),
        ("DataBus",        ACCENT),
        ("Memory",         STEEL),
        ("Vault",          STEEL),
        ("Monitoring",     STEEL),
    ]
    for i, (label, color) in enumerate(components):
        bx = diag_l + i * (box_w + gap)
        r = rect(slide, bx, diag_t, box_w, box_h, fill=color)
        txbox(slide, label, bx + Inches(0.05), diag_t + Inches(0.1),
              box_w - Inches(0.1), box_h - Inches(0.15),
              size=11, bold=(label == "DataBus"), color=WHITE, align=PP_ALIGN.CENTER)

    # Arrows between boxes (simple text "→")
    for i in range(len(components) - 1):
        ax = diag_l + i * (box_w + gap) + box_w + Inches(0.01)
        txbox(slide, "→", ax, diag_t + Inches(0.12), gap + Inches(0.1), Inches(0.35),
              size=16, color=RGBColor(0xAE,0xBF,0xC8), align=PP_ALIGN.CENTER)

    # Caption
    txbox(slide, "All 18 Aegis OS components route exclusively through the DataBus",
          Inches(0.5), Inches(5.0), Inches(12.5), Inches(0.4),
          size=13, color=RGBColor(0xAE,0xBF,0xC8), italic=True, align=PP_ALIGN.LEFT)

    txbox(slide, "CS-686 | Aegis DataBus | Meron Shibiru | March 9, 2026",
          Inches(0.3), SLIDE_H - Inches(0.40), SLIDE_W - Inches(0.6), Inches(0.38),
          size=9, color=RGBColor(0x5D,0x6D,0x7E), align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Welcome. Today I'm presenting the Aegis DataBus — the central nervous system of Aegis OS.\n"
        "Every one of the 18 OS components communicates exclusively through this component.\n"
        "If the DataBus dies, the OS stops. It is first to boot and last to shut down.\n"
        "I'll walk through what it is, how it works, the design patterns and security model, "
        "and show that every EDD requirement has a passing test.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is Aegis OS?
# ═══════════════════════════════════════════════════════════════════════════════
def slide_aegis_os(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "What is Aegis OS?", "Context: Why the DataBus matters")
    add_footer(slide)

    bullets = [
        "AI-Native Operating System for autonomous AI agents",
        "Runs agents in isolated microVMs on bare-metal Linux",
        "18 components — every one depends on the DataBus",
        "DataBus: first to boot, last to shut down",
        "If the DataBus dies → the OS stops (Severity-1 incident)",
    ]
    bullet_block(slide, bullets, t=Inches(1.65), item_h=Inches(0.72), size=18)

    # Callout box
    r = rect(slide, Inches(7.8), Inches(2.0), Inches(5.1), Inches(2.5),
             fill=NAVY, line=ACCENT)
    txbox(slide, "Zero-dependency root\n\nMust start FIRST.\nAll 18 components wait\nfor the DataBus to be ready.",
          Inches(7.95), Inches(2.1), Inches(4.9), Inches(2.3),
          size=15, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Aegis OS is an AI-native OS — it schedules and runs autonomous AI agents in microVMs.\n"
        "18 distinct components handle tasks like routing, planning, memory, vaulting, monitoring.\n"
        "None of them talk directly to each other — everything goes through the DataBus.\n"
        "This makes the DataBus a zero-dependency root: it boots first, shuts down last.\n"
        "A DataBus failure is a P1 system-wide incident.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What is the DataBus?
# ═══════════════════════════════════════════════════════════════════════════════
def slide_what_is(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "What is the DataBus?", "The ONLY inter-component communication channel")
    add_footer(slide)

    bullets = [
        "Central nervous system — the ONLY way components communicate",
        "NO direct function calls, NO shared memory, NO REST between components",
        "Everything goes through publish/subscribe messaging",
    ]
    bullet_block(slide, bullets, t=Inches(1.7), item_h=Inches(0.65), size=18)

    # Tech stack bar
    rect(slide, Inches(0.5), Inches(3.3), Inches(12.3), Inches(0.55), fill=NAVY)
    txbox(slide, "Technology Stack:   Go  ·  NATS JetStream 2.10+  ·  MemoryClient Interface  ·  NKey Auth  ·  Docker",
          Inches(0.6), Inches(3.32), Inches(12.1), Inches(0.50),
          size=15, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    # Three pillars
    pillars = [
        ("Durable Streams", "7-day retention\n10 GB max\n3 replicas"),
        ("Zero-Loss Delivery", "Outbox pattern\nExplicit ACK\nExponential backoff"),
        ("Security First", "NKey auth\nSubject ACLs\nTLS 1.3"),
    ]
    for i, (title, body) in enumerate(pillars):
        bx = Inches(0.5) + i * Inches(4.15)
        rect(slide, bx, Inches(4.1), Inches(3.9), Inches(2.7), fill=STEEL)
        txbox(slide, title, bx + Inches(0.15), Inches(4.2), Inches(3.6), Inches(0.55),
              size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        txbox(slide, body, bx + Inches(0.15), Inches(4.75), Inches(3.6), Inches(1.9),
              size=14, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "The DataBus is the sole communication channel — no component ever calls another directly.\n"
        "Built on NATS JetStream which gives us durable, persistent, at-least-once delivery.\n"
        "Go for the implementation. NKey auth for Zero Trust. Docker for deployment.\n"
        "The MemoryClient interface allows us to swap between a mock and HTTP implementation.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What the DataBus is NOT
# ═══════════════════════════════════════════════════════════════════════════════
def slide_not(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "What the DataBus is NOT", "Explicit design boundaries")
    add_footer(slide)

    not_items = [
        ("NOT an API gateway",   "No HTTP routing, no load-balancing decisions"),
        ("NOT a database",       "Messages have TTL; it is not a persistent store"),
        ("NOT a service mesh",   "No Istio/Envoy sidecar — messaging only"),
        ("No business logic",    "DataBus does not know what agents do or mean"),
        ("No agent awareness",   "Agnostic to payload semantics — routes by subject"),
    ]

    for i, (title, detail) in enumerate(not_items):
        by = Inches(1.7) + i * Inches(0.95)
        rect(slide, Inches(0.5), by, Inches(3.5), Inches(0.75), fill=RGBColor(0xC0,0x39,0x2B))
        txbox(slide, title, Inches(0.6), by + Inches(0.1), Inches(3.3), Inches(0.6),
              size=15, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
        txbox(slide, detail, Inches(4.2), by + Inches(0.12), Inches(8.5), Inches(0.55),
              size=15, color=DARK_GRAY)

    add_speaker_notes(slide,
        "It's important to understand what the DataBus doesn't do — this keeps it lean and testable.\n"
        "It routes messages by subject pattern; it doesn't inspect or transform payloads.\n"
        "It's not a replacement for a database — messages have configurable TTL.\n"
        "Business logic lives in the components, not in the bus.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Core Responsibilities
# ═══════════════════════════════════════════════════════════════════════════════
def slide_responsibilities(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Core Responsibilities", "Six guarantees the DataBus provides")
    add_footer(slide)

    items = [
        ("1", "Fan-out messages to matching subscribers",  "Broadcast / wildcard pub-sub (FR-DB-001, FR-DB-005)"),
        ("2", "Persist to durable JetStream streams",      "Survive restarts; replay from any offset (FR-DB-003, FR-DB-008)"),
        ("3", "Zero message loss — Outbox Pattern",        "Atomic DB write + async relay (FR-DB-011)"),
        ("4", "Enforce security",                          "NKey auth + subject-level ACLs (SR-DB-001, SR-DB-006)"),
        ("5", "Health & metrics for Monitoring",           "/metrics, /healthz, /varz, /connz, /jsz (Interface 4)"),
        ("6", "DEGRADED-HOLD mode",                        "Never crash when dependencies are down — use MockMemoryClient"),
    ]

    for i, (num, title, detail) in enumerate(items):
        by = Inches(1.65) + i * Inches(0.88)
        rect(slide, Inches(0.5), by, Inches(0.55), Inches(0.65), fill=ACCENT)
        txbox(slide, num, Inches(0.5), by + Inches(0.08), Inches(0.55), Inches(0.5),
              size=20, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        rect(slide, Inches(1.15), by, Inches(11.65), Inches(0.65), fill=WHITE,
             line=RGBColor(0xBD,0xC3,0xC7))
        txbox(slide, title, Inches(1.25), by + Inches(0.04), Inches(4.8), Inches(0.35),
              size=15, bold=True, color=STEEL)
        txbox(slide, detail, Inches(6.1), by + Inches(0.14), Inches(6.5), Inches(0.42),
              size=12, color=DARK_GRAY, italic=True)

    add_speaker_notes(slide,
        "Six core responsibilities cover the full lifecycle of a message in the system.\n"
        "Fan-out ensures all subscribers get events. Persistence enables replay after failure.\n"
        "The Outbox pattern is the key zero-loss guarantee — I'll explain it in depth shortly.\n"
        "DEGRADED-HOLD is critical: if Memory is down, we switch to MockMemoryClient and keep running.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — The 6 JetStream Streams
# ═══════════════════════════════════════════════════════════════════════════════
def slide_streams(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "The 6 JetStream Streams (Channels)", "Domain-scoped durable message logs")
    add_footer(slide)

    headers = ["Stream", "Subject Pattern", "Key Publishers", "Key Subscribers"]
    rows = [
        ["AEGIS_TASKS",      "aegis.tasks.>",      "Task Router, Planner",       "Task Planner, Agent Mgr"],
        ["AEGIS_AGENTS",     "aegis.agents.>",     "Agent Manager",              "Self-Healing, Memory Mgr"],
        ["AEGIS_RUNTIME",    "aegis.runtime.>",    "Runtime Abstraction",        "Agent Mgr, Monitoring"],
        ["AEGIS_VAULT",      "aegis.vault.>",      "Permission Manager",         "Runtime Abstraction"],
        ["AEGIS_MEMORY",     "aegis.memory.>",     "Memory & Context Mgr",       "Self-Healing, Knowledge"],
        ["AEGIS_MONITORING", "aegis.monitoring.>", "Monitoring & Observability", "All components"],
    ]
    two_col_table(slide, headers, rows,
                  l=Inches(0.4), t=Inches(1.65),
                  w=Inches(12.5), h=Inches(4.2))

    # Config badge
    rect(slide, Inches(0.4), Inches(6.2), Inches(12.5), Inches(0.7), fill=NAVY)
    txbox(slide,
          "Config (all streams):  7-day retention  ·  10 GB max  ·  3 replicas  ·  Explicit ACK  ·  120s dedup window",
          Inches(0.55), Inches(6.25), Inches(12.2), Inches(0.6),
          size=13, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Six streams, each scoping a logical domain of Aegis OS.\n"
        "Subject patterns use '>' wildcard — e.g. aegis.tasks.> captures all task events.\n"
        "Three replicas across the cluster give us fault tolerance — any two nodes can fail.\n"
        "120s dedup window catches double-publishes within the window using the CloudEvents id field.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CloudEvents Message Format
# ═══════════════════════════════════════════════════════════════════════════════
def slide_cloudevents(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "CloudEvents 1.0 Message Format", "Standard envelope — FR-DB-010 schema validation")
    add_footer(slide)

    # Left: field list
    fields = [
        ("specversion",     '"1.0"',                          "Fixed — identifies CloudEvents version"),
        ("id",              "uuid-v4",                        "Used for deduplication (120s window)"),
        ("source",          '"aegis/task-router"',            "Originating component"),
        ("type",            '"aegis.tasks.routed"',           "Event type — matches NATS subject"),
        ("time",            "RFC3339 timestamp",              "Event creation time"),
        ("correlationid",   "uuid (trace chain)",             "Links related events across components"),
        ("datacontenttype", '"application/json"',             "Always JSON"),
        ("data",            '{ ... }',                        "Actual payload — never logged"),
    ]

    rect(slide, Inches(0.4), Inches(1.6), Inches(7.8), Inches(5.5), fill=WHITE,
         line=RGBColor(0xBD,0xC3,0xC7))

    txbox(slide, "Field", Inches(0.5), Inches(1.65), Inches(2.1), Inches(0.4),
          size=13, bold=True, color=NAVY)
    txbox(slide, "Example Value", Inches(2.6), Inches(1.65), Inches(2.5), Inches(0.4),
          size=13, bold=True, color=NAVY)
    txbox(slide, "Purpose", Inches(5.1), Inches(1.65), Inches(2.9), Inches(0.4),
          size=13, bold=True, color=NAVY)
    rect(slide, Inches(0.4), Inches(2.05), Inches(7.8), Pt(2), fill=ACCENT)

    for i, (field, val, purpose) in enumerate(fields):
        by = Inches(2.1) + i * Inches(0.55)
        fill = WHITE if i % 2 == 0 else RGBColor(0xD6,0xEA,0xF8)
        rect(slide, Inches(0.4), by, Inches(7.8), Inches(0.55), fill=fill)
        txbox(slide, field,   Inches(0.5), by + Pt(3), Inches(2.1), Inches(0.5), size=12, bold=True,  color=STEEL)
        txbox(slide, val,     Inches(2.6), by + Pt(3), Inches(2.5), Inches(0.5), size=11, color=RGBColor(0x1A,0x5,0x76))
        txbox(slide, purpose, Inches(5.1), by + Pt(3), Inches(2.9), Inches(0.5), size=11, color=DARK_GRAY)

    # Right: Enforcement callout
    rect(slide, Inches(8.5), Inches(1.6), Inches(4.4), Inches(3.5), fill=NAVY)
    txbox(slide, "Enforcement",
          Inches(8.65), Inches(1.7), Inches(4.1), Inches(0.5),
          size=16, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    enforcement = (
        "▸ PublishValidated()\n"
        "▸ PublishWithACL()\n"
        "▸ PublishAsync()\n"
        "▸ PublishBatch()\n\n"
        "All call envelope.Validate()\nbefore sending to NATS.\n\n"
        "Invalid → error returned\n+ validation_errors_total++"
    )
    txbox(slide, enforcement, Inches(8.65), Inches(2.25), Inches(4.1), Inches(2.7),
          size=13, color=WHITE)

    rect(slide, Inches(8.5), Inches(5.3), Inches(4.4), Inches(1.5), fill=RGBColor(0xC0,0x39,0x2B))
    txbox(slide, "Security Note\n\ndata field is NEVER written\nto logs or audit trail\n(SR-DB-005)",
          Inches(8.65), Inches(5.35), Inches(4.1), Inches(1.4),
          size=13, bold=False, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "All messages must be valid CloudEvents 1.0 — this is enforced at every publish call.\n"
        "The 'id' field is used for 120-second deduplication in JetStream — same id twice is dropped.\n"
        "correlationid chains events across components for distributed tracing.\n"
        "Critical: the 'data' field is never logged. Audit logs only store metadata: subject, size, correlationid.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Design Patterns
# ═══════════════════════════════════════════════════════════════════════════════
def slide_patterns(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Design Patterns", "Proven patterns for reliability and scale")
    add_footer(slide)

    patterns = [
        ("Publish-Subscribe",         "Decouple producers from consumers",             "FR-DB-001, FR-DB-005"),
        ("Outbox Pattern",            "Zero-loss guarantee — atomic DB + async relay",  "FR-DB-011"),
        ("Durable Consumer",          "Survive restarts; resume from last ACK'd seq",   "FR-DB-003"),
        ("Dead Letter Queue (DLQ)",   "Isolate failed messages after 5 NAKs",           "FR-DB-009, SR-DB-006"),
        ("Queue Groups",              "Competing consumers — horizontal scaling",        "FR-DB-004"),
        ("Adapter / Dep. Inversion",  "MemoryClient interface: swap Mock ↔ HTTP",        "Testability"),
    ]

    cols = [
        ("Pattern",     Inches(0.5),  Inches(3.2)),
        ("Purpose",     Inches(3.8),  Inches(5.2)),
        ("Requirement", Inches(9.1),  Inches(3.8)),
    ]

    # Header row
    for label, lx, cw in cols:
        rect(slide, lx, Inches(1.63), cw, Inches(0.5), fill=NAVY)
        txbox(slide, label, lx + Inches(0.1), Inches(1.65), cw - Inches(0.15), Inches(0.45),
              size=13, bold=True, color=WHITE)

    for i, (pat, purpose, req) in enumerate(patterns):
        by = Inches(2.18) + i * Inches(0.8)
        fill = WHITE if i % 2 == 0 else RGBColor(0xD6,0xEA,0xF8)
        for lx, cw in [(Inches(0.5), Inches(3.2)), (Inches(3.8), Inches(5.2)), (Inches(9.1), Inches(3.8))]:
            rect(slide, lx, by, cw, Inches(0.78), fill=fill, line=RGBColor(0xBD,0xC3,0xC7))
        txbox(slide, pat,     Inches(0.6),  by + Pt(6), Inches(3.0), Inches(0.65), size=13, bold=True,  color=STEEL)
        txbox(slide, purpose, Inches(3.9),  by + Pt(6), Inches(5.0), Inches(0.65), size=13, color=DARK_GRAY)
        txbox(slide, req,     Inches(9.2),  by + Pt(6), Inches(3.6), Inches(0.65), size=12, color=MID_BLUE, italic=True)

    add_speaker_notes(slide,
        "Six patterns working together to deliver the guarantees in the EDD.\n"
        "Pub-Sub is the foundation. Outbox is the most important for correctness.\n"
        "DLQ prevents poison messages from blocking streams.\n"
        "The Adapter pattern is key for testability — we run tests without the Memory HTTP service.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Key Design Decisions
# ═══════════════════════════════════════════════════════════════════════════════
def slide_decisions(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Key Design Decisions", "Why we built it this way")
    add_footer(slide)

    decisions = [
        (
            "Schema validation at the API boundary",
            "PublishValidated / PublishWithACL enforce CloudEvents.\n"
            "NATS itself never inspects payloads — validation is our responsibility.",
            "Centralised — one place to add rules"
        ),
        (
            "ACL in application layer",
            "NATS auth is connection-level (NKey).\n"
            "Subject-level rules live in pkg/security/acl.go — CheckPublish / CheckSubscribe.",
            "Fine-grained without NATS operator config"
        ),
        (
            "Fallback to Mock on Memory failure",
            "FallbackClient wraps HTTP client. On Ping() failure → switches to MockMemoryClient.\n"
            "DataBus stays up. Outbox uses in-memory mock. DEGRADED-HOLD — never crash.",
            "Resilience without data loss"
        ),
        (
            "Stream setup with retries",
            "10× retries, 2s delay for 3-node Raft quorum to form.\n"
            "Prevents race condition at startup where JetStream isn't ready yet.",
            "Startup reliability"
        ),
    ]

    for i, (title, body, outcome) in enumerate(decisions):
        by = Inches(1.65) + i * Inches(1.35)
        rect(slide, Inches(0.4), by, Inches(12.5), Inches(1.25), fill=WHITE,
             line=RGBColor(0xBD,0xC3,0xC7))
        rect(slide, Inches(0.4), by, Inches(0.12), Inches(1.25), fill=ACCENT)
        txbox(slide, title, Inches(0.65), by + Inches(0.05), Inches(9.0), Inches(0.4),
              size=15, bold=True, color=NAVY)
        txbox(slide, body, Inches(0.65), by + Inches(0.42), Inches(9.0), Inches(0.75),
              size=12, color=DARK_GRAY)
        # Outcome badge
        rect(slide, Inches(9.8), by + Inches(0.3), Inches(3.0), Inches(0.6), fill=STEEL)
        txbox(slide, outcome, Inches(9.9), by + Inches(0.33), Inches(2.85), Inches(0.55),
              size=11, color=WHITE, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "Four decisions that shaped the implementation.\n"
        "Validation at the boundary means downstream consumers never see malformed messages.\n"
        "Application-layer ACLs give us subject-level control without complex NATS operator setup.\n"
        "The fallback pattern means Memory being down doesn't crash the DataBus — critical for HA.\n"
        "The retry loop on stream setup handles the race where Raft hasn't formed a quorum yet.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Security
# ═══════════════════════════════════════════════════════════════════════════════
def slide_security(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Security Model", "Zero Trust — never trust, always verify")
    add_footer(slide)

    controls = [
        ("No payload in logs",   "SR-DB-005",
         "Audit stores: subject · size · correlationId · traceID only. NEVER payload data."),
        ("Subject ACLs",         "SR-DB-006",
         "CheckPublish() / CheckSubscribe() on every call. DLQ admin-only."),
        ("NKey Authentication",  "SR-DB-001, SR-DB-004",
         "Per-component Ed25519 NKey. OpenBao (production) or env var fallback."),
        ("TLS 1.3",              "SR-DB-002",
         "make up-tls — mTLS between all nodes and clients. Plaintext rejected."),
        ("Plaintext rejected",   "SR-DB-001",
         "NKey-enabled NATS drops unauthenticated connections."),
    ]

    for i, (control, req_id, detail) in enumerate(controls):
        by = Inches(1.65) + i * Inches(0.98)
        rect(slide, Inches(0.4), by, Inches(12.5), Inches(0.88), fill=NAVY)
        # req badge
        rect(slide, Inches(0.4), by, Inches(1.7), Inches(0.88), fill=ACCENT)
        txbox(slide, req_id, Inches(0.42), by + Inches(0.15), Inches(1.65), Inches(0.6),
              size=11, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txbox(slide, control, Inches(2.25), by + Inches(0.05), Inches(3.5), Inches(0.4),
              size=14, bold=True, color=WHITE)
        txbox(slide, detail, Inches(2.25), by + Inches(0.45), Inches(10.5), Inches(0.4),
              size=12, color=RGBColor(0xAE,0xBF,0xC8))

    add_speaker_notes(slide,
        "Security is defence-in-depth across multiple layers.\n"
        "Audit logs never contain payload — only metadata. This is SR-DB-005.\n"
        "ACLs are enforced in the application because NATS auth is connection-level only.\n"
        "NKeys are Ed25519 keypairs — when OpenBao is running, seeds come from the vault.\n"
        "In demo mode, seeds come from environment variables — the code path is the same.\n"
        "TLS 1.3 is available with 'make up-tls' and generated self-signed certs.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — High Availability
# ═══════════════════════════════════════════════════════════════════════════════
def slide_ha(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "High Availability", "Designed to survive failures at every layer")
    add_footer(slide)

    items = [
        ("3-node NATS cluster",   "Raft consensus · leader election < 5s (FR-DB-007)"),
        ("Fallback MemoryClient", "DEGRADED-HOLD when Memory service is down — MockMemoryClient"),
        ("Exponential backoff",   "Stream setup retries · Outbox relay retries — never tight-loop"),
        ("Heartbeat",             "aegis.health.databus — Self-Healing Controller monitors liveness"),
        ("Auto-reconnect",        "MaxReconnects(-1) · exponential backoff on client disconnect"),
    ]

    for i, (title, detail) in enumerate(items):
        by = Inches(1.65) + i * Inches(0.9)
        # Left colored bar
        rect(slide, Inches(0.4), by, Inches(4.0), Inches(0.78), fill=STEEL)
        txbox(slide, title, Inches(0.55), by + Inches(0.17), Inches(3.7), Inches(0.5),
              size=15, bold=True, color=WHITE)
        # Right detail
        rect(slide, Inches(4.5), by, Inches(8.4), Inches(0.78), fill=WHITE,
             line=RGBColor(0xBD,0xC3,0xC7))
        txbox(slide, detail, Inches(4.65), by + Inches(0.17), Inches(8.1), Inches(0.5),
              size=14, color=DARK_GRAY)

    # Recovery sequence
    rect(slide, Inches(0.4), Inches(6.15), Inches(12.5), Inches(0.75), fill=NAVY)
    txbox(slide, "Failure Recovery:  Heartbeat timeout → Raft elects leader < 5s → clients reconnect → replay durable sequences → RecoveryCompleted",
          Inches(0.55), Inches(6.2), Inches(12.2), Inches(0.65),
          size=12, color=ACCENT, italic=True, align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "High availability at every layer — the cluster, the client, and the dependencies.\n"
        "3-node Raft gives us N+1 fault tolerance — one node can die without service interruption.\n"
        "Leader election completes in under 5 seconds — verified by TestFRDB007_ClusterFailover.\n"
        "The heartbeat on aegis.health.databus is what the Self-Healing Controller watches.\n"
        "DEGRADED-HOLD is unique: if Memory HTTP is unreachable, we fall back to an in-process mock.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Test Harness
# ═══════════════════════════════════════════════════════════════════════════════
def slide_tests(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Test Harness — All Requirements Verified", "Every EDD requirement has a passing test")
    add_footer(slide)

    headers = ["Test ID", "Requirement", "What It Verifies"]
    rows = [
        ["TC001 / TC001b", "FR-DB-001",         "Pub/sub latency < 5ms P99"],
        ["TC002",          "FR-DB-004",         "Queue group — each msg delivered to exactly 1 subscriber"],
        ["TC003",          "FR-DB-003",         "Durable consumer recovery across reconnect"],
        ["TC004",          "FR-DB-008",         "Outbox relay replay after simulated crash"],
        ["TC005",          "FR-DB-009, SR-DB-006", "DLQ after 5 NAKs; non-admin subscribe denied"],
        ["TC006",          "FR-DB-011",         "Outbox zero-loss — message delivered on relay start"],
        ["FR-DB-002",      "FR-DB-002",         "Request-reply — reply within 5s timeout"],
        ["FR-DB-006",      "FR-DB-006",         "Priority — health subjects before resource subjects"],
        ["Benchmark",      "NFR-DB-001/005",    "≥ 50K msg/s throughput; 5ms P99 latency"],
    ]
    two_col_table(slide, headers, rows,
                  l=Inches(0.4), t=Inches(1.63),
                  w=Inches(12.5), h=Inches(5.2))

    add_speaker_notes(slide,
        "Nine test cases covering every functional and security requirement in the EDD.\n"
        "TC001 uses a real NATS server — it measures actual round-trip latency.\n"
        "TC005 is particularly important: it verifies both the DLQ routing AND the ACL denial for non-admin.\n"
        "The benchmark runs 50K messages and measures P99 latency — verifies NFR-DB-001.\n"
        "Run with: make test-integration (integration) or make test (all packages).")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Demo Architecture
# ═══════════════════════════════════════════════════════════════════════════════
def slide_demo_arch(prs):
    slide = add_slide(prs)
    body_bg(slide)
    header_bar(slide, "Demo Architecture — 6 Components", "Simulating the full Aegis OS message flow")
    add_footer(slide)

    components = [
        ("I/O Layer",       "Publishes user actions\naegis.tasks.submitted",    MID_BLUE),
        ("Orchestrator",    "Task Router + Planner\n+ Agent Manager (queue grp)", STEEL),
        ("Memory",          "Context Manager\nRequest-reply responder",           RGBColor(0x17,0x7E,0x5E)),
        ("Vault",           "Permission Manager\naegis.vault.>",                  RGBColor(0x7D,0x3C,0x98)),
        ("Agent",           "Agent runtime\nCompletes tasks",                     RGBColor(0xC0,0x39,0x2B)),
        ("Monitoring",      "Subscribes aegis.>\n(wildcard — all events)",        RGBColor(0xD3,0x54,0x00)),
    ]

    box_w = Inches(1.95)
    box_h = Inches(1.5)
    row_top_y = Inches(2.0)
    row_bot_y = Inches(3.9)
    start_x = Inches(0.5)
    gap = Inches(0.25)

    for i, (name, desc, color) in enumerate(components):
        if i < 3:
            bx = start_x + i * (box_w + gap)
            by = row_top_y
        else:
            bx = start_x + (i - 3) * (box_w + gap)
            by = row_bot_y
        rect(slide, bx, by, box_w, box_h, fill=color)
        txbox(slide, name, bx + Inches(0.08), by + Inches(0.08),
              box_w - Inches(0.15), Inches(0.45),
              size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        txbox(slide, desc, bx + Inches(0.08), by + Inches(0.55),
              box_w - Inches(0.15), Inches(0.85),
              size=11, color=WHITE, align=PP_ALIGN.CENTER)

    # DataBus center column
    db_x = Inches(6.7)
    db_y = Inches(1.9)
    db_w = Inches(2.4)
    db_h = Inches(3.6)
    rect(slide, db_x, db_y, db_w, db_h, fill=NAVY)
    txbox(slide, "DataBus\n\nNATS JetStream\n3-node cluster\n\n/metrics  /healthz\n/varz  /connz  /jsz",
          db_x + Inches(0.1), db_y + Inches(0.2),
          db_w - Inches(0.2), db_h - Inches(0.3),
          size=13, bold=False, color=ACCENT, align=PP_ALIGN.CENTER)

    # Arrows
    for ax, ay in [(Inches(6.55), Inches(2.4)), (Inches(6.55), Inches(4.2))]:
        txbox(slide, "↔", ax, ay, Inches(0.3), Inches(0.4), size=18, color=ACCENT,
              align=PP_ALIGN.CENTER)

    # Right: demo command
    rect(slide, Inches(9.3), Inches(1.9), Inches(3.6), Inches(3.5), fill=DARK_GRAY)
    txbox(slide, "Run the Demo",
          Inches(9.45), Inches(1.95), Inches(3.3), Inches(0.45),
          size=14, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
    demo_cmds = (
        "make up\n"
        "./bin/aegis-databus &\n"
        "./bin/aegis-demo\n\n"
        "# Or one command:\n"
        "make demo-full\n\n"
        "Grafana:  :3000\n"
        "NATS UI:  :8222/jsz"
    )
    txbox(slide, demo_cmds, Inches(9.45), Inches(2.45), Inches(3.3), Inches(2.85),
          size=12, color=WHITE)

    add_speaker_notes(slide,
        "The demo uses 6 simulated components that mirror real Aegis OS components.\n"
        "I/O Layer publishes user actions. Orchestrator uses queue group for horizontal scaling.\n"
        "Memory uses request-reply — Orchestrator asks for preferences, Memory responds.\n"
        "Monitoring subscribes to aegis.> wildcard — it sees ALL events across all streams.\n"
        "The DataBus is the only entity in the middle — everything routes through it.\n"
        "Run 'make demo-full' for one-command orchestrated demo with Grafana dashboard.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Summary & Takeaways
# ═══════════════════════════════════════════════════════════════════════════════
def slide_summary(prs):
    slide = add_slide(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=NAVY)
    rect(slide, 0, 0, Inches(0.12), SLIDE_H, fill=ACCENT)
    rect(slide, Inches(0.3), Inches(1.5), Inches(12.7), Pt(3), fill=ACCENT)

    txbox(slide, "Summary & Takeaways",
          Inches(0.5), Inches(0.4), Inches(12.5), Inches(0.9),
          size=40, bold=True, color=WHITE)

    takeaways = [
        ("Central backbone",    "The DataBus is the ONLY inter-component channel in Aegis OS"),
        ("Zero-loss",           "Outbox pattern + durable JetStream = no message loss across crashes"),
        ("Secure",              "NKey auth + TLS 1.3 + subject ACLs + audit-without-payload"),
        ("Resilient",           "3-node Raft + DEGRADED-HOLD + exponential backoff everywhere"),
        ("Fully tested",        "TC001–TC006, FR-DB-002, FR-DB-006 — every EDD requirement covered"),
        ("Live demo ready",     "Zero Trust, DLQ, cluster failover, outbox, request-reply — all runnable"),
    ]

    for i, (key, val) in enumerate(takeaways):
        by = Inches(1.7) + i * Inches(0.85)
        rect(slide, Inches(0.5), by, Inches(2.8), Inches(0.68), fill=ACCENT)
        txbox(slide, key, Inches(0.6), by + Inches(0.12), Inches(2.6), Inches(0.5),
              size=15, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txbox(slide, val, Inches(3.5), by + Inches(0.12), Inches(9.5), Inches(0.5),
              size=15, color=WHITE)

    txbox(slide, "Questions?  —  github.com/mlim3/cerberOS  ·  CS-686 | March 9, 2026",
          Inches(0.5), SLIDE_H - Inches(0.55), Inches(12.5), Inches(0.45),
          size=12, color=RGBColor(0x5D,0x6D,0x7E), align=PP_ALIGN.CENTER)

    add_speaker_notes(slide,
        "To summarise:\n"
        "The Aegis DataBus implements the central message transport for all 18 OS components.\n"
        "Key achievements: zero message loss via Outbox, Zero Trust via NKey/OpenBao, "
        "3-node Raft cluster with < 5s failover, and DEGRADED-HOLD resilience.\n"
        "Every functional, NFR, and security requirement in the EDD has a passing automated test.\n"
        "I'm happy to demo any of the 12 demo scenarios live.")
    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════
def main():
    prs = new_prs()

    slide_title(prs)
    slide_aegis_os(prs)
    slide_what_is(prs)
    slide_not(prs)
    slide_responsibilities(prs)
    slide_streams(prs)
    slide_cloudevents(prs)
    slide_patterns(prs)
    slide_decisions(prs)
    slide_security(prs)
    slide_ha(prs)
    slide_tests(prs)
    slide_demo_arch(prs)
    slide_summary(prs)

    out = "/Users/merakera/Spring-2026/cs-686/cerberOS/Aegis-DataBus-Presentation.pptx"
    prs.save(out)
    print(f"Saved: {out}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
